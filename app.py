"""
Consulting Pitch Deck Generator
Generates a downloadable .pptx file structured like a top-tier consulting deck,
based on a client situation, problem description, and chosen analytical framework.

Customisations beyond starter:
- Research context field: lets users paste real client research to ground the output
- Framework picker: SWOT, Porter's Five Forces, McKinsey 7S, Three Horizons, MECE issue tree
- Polished PowerPoint design: header strip, page numbers, footer, refined typography
- Options comparison diagram: replaces bullet list with a 3-column comparison table
- 8 visual theme palettes auto-suggested by industry

Built by Gan Kai Li as part of preparing for Big Four consulting work.
"""

import streamlit as st
from google import genai
from dotenv import load_dotenv
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import os
import json

# ---- Configuration ----
env_path = Path(__file__).parent / ".env"
load_dotenv(dotenv_path=env_path)

try:
    API_KEY = st.secrets["GOOGLE_API_KEY"]
except (KeyError, FileNotFoundError):
    API_KEY = os.getenv("GOOGLE_API_KEY")

st.set_page_config(page_title="Consulting Pitch Deck Generator", page_icon="📊", layout="wide")
st.title("Consulting Pitch Deck Generator")
st.caption("Generate structured pitch decks in the style of top-tier consulting firms. Built for aspiring and current consultants.")

if not API_KEY:
    st.error("No API key found. Check your .env file contains GOOGLE_API_KEY=your-key")
    st.stop()

client = genai.Client(api_key=API_KEY)


# ---- Framework definitions ----
FRAMEWORKS = {
    "SCR (Situation-Complication-Resolution)": {
        "description": "Classic McKinsey storyline. Best for general strategy problems.",
        "sections": ["Executive Summary", "Situation", "Complication", "Resolution Options", "Recommendation"],
    },
    "Porter's Five Forces": {
        "description": "Industry attractiveness analysis. Best for market entry or competitive strategy.",
        "sections": ["Executive Summary", "Industry Rivalry", "Threat of New Entrants", "Buyer & Supplier Power", "Threat of Substitutes", "Strategic Implications"],
    },
    "SWOT Analysis": {
        "description": "Strengths, Weaknesses, Opportunities, Threats. Best for situation assessment.",
        "sections": ["Executive Summary", "Strengths", "Weaknesses", "Opportunities", "Threats", "Strategic Priorities"],
    },
    "McKinsey 7S": {
        "description": "Organisational design framework. Best for operating model or change problems.",
        "sections": ["Executive Summary", "Strategy & Structure", "Systems & Style", "Staff & Skills", "Shared Values", "Recommendations"],
    },
    "Three Horizons": {
        "description": "Growth strategy across time. Best for portfolio or transformation decisions.",
        "sections": ["Executive Summary", "Horizon 1 — Core Business", "Horizon 2 — Adjacent Growth", "Horizon 3 — Future Bets", "Capital Allocation", "Recommendation"],
    },
    "MECE Issue Tree": {
        "description": "Structured problem decomposition. Best for root-cause analysis.",
        "sections": ["Executive Summary", "Problem Statement", "Branch 1 Analysis", "Branch 2 Analysis", "Branch 3 Analysis", "Synthesis & Recommendation"],
    },
}


# ---- Visual theme palettes ----
PALETTES = {
    "Corporate Navy": {
        "primary": RGBColor(0x0A, 0x2A, 0x4A),
        "accent":  RGBColor(0xC0, 0x39, 0x2B),
        "text":    RGBColor(0x33, 0x33, 0x33),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xF5, 0xF5, 0xF5),
    },
    "Banking Blue": {
        "primary": RGBColor(0x1F, 0x3A, 0x5F),
        "accent":  RGBColor(0xE8, 0xB5, 0x4B),
        "text":    RGBColor(0x2C, 0x2C, 0x2C),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xF2, 0xF5, 0xF8),
    },
    "Energy Bold": {
        "primary": RGBColor(0x1A, 0x1A, 0x1A),
        "accent":  RGBColor(0xE6, 0x4A, 0x19),
        "text":    RGBColor(0x33, 0x33, 0x33),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xF7, 0xF3, 0xEE),
    },
    "Tech Modern": {
        "primary": RGBColor(0x1A, 0x2E, 0x4F),
        "accent":  RGBColor(0x00, 0xB8, 0xA9),
        "text":    RGBColor(0x33, 0x33, 0x33),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xF0, 0xF8, 0xF7),
    },
    "Consumer Vibrant": {
        "primary": RGBColor(0x6B, 0x1F, 0x5C),
        "accent":  RGBColor(0xFF, 0x6B, 0x35),
        "text":    RGBColor(0x33, 0x33, 0x33),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xFB, 0xF2, 0xF6),
    },
    "Healthcare Calm": {
        "primary": RGBColor(0x1B, 0x4D, 0x5A),
        "accent":  RGBColor(0x4A, 0xA8, 0x8E),
        "text":    RGBColor(0x33, 0x33, 0x33),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xEF, 0xF5, 0xF4),
    },
    "Government Formal": {
        "primary": RGBColor(0x2F, 0x2F, 0x2F),
        "accent":  RGBColor(0x8B, 0x1A, 0x1A),
        "text":    RGBColor(0x2C, 0x2C, 0x2C),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xF4, 0xF4, 0xF4),
    },
    "Industrial Strong": {
        "primary": RGBColor(0x2D, 0x3E, 0x50),
        "accent":  RGBColor(0xF3, 0x9C, 0x12),
        "text":    RGBColor(0x33, 0x33, 0x33),
        "footer":  RGBColor(0x88, 0x88, 0x88),
        "bg_box":  RGBColor(0xF5, 0xF1, 0xEA),
    },
}


def suggest_palette(client_name, industry):
    """Suggest a palette based on simple keyword matching against client and industry."""
    text = (client_name + " " + industry).lower()
    if any(kw in text for kw in ["bank", "financ", "insurance", "wealth", "khazanah"]):
        return "Banking Blue"
    if any(kw in text for kw in ["oil", "gas", "petronas", "energy", "utility", "utilities", "power"]):
        return "Energy Bold"
    if any(kw in text for kw in ["tech", "software", "saas", "digital", "ai ", "fintech", "platform"]):
        return "Tech Modern"
    if any(kw in text for kw in ["retail", "consumer", "fmcg", "fashion", "hospitality", "food", "airline", "airasia", "aviation"]):
        return "Consumer Vibrant"
    if any(kw in text for kw in ["health", "pharma", "medical", "hospital", "biotech", "clinic"]):
        return "Healthcare Calm"
    if any(kw in text for kw in ["govern", "ministry", "sovereign", "public sector", "ngo", "nonprofit"]):
        return "Government Formal"
    if any(kw in text for kw in ["manufactur", "logistics", "construction", "industrial", "shipping", "mining"]):
        return "Industrial Strong"
    return "Corporate Navy"


# ---- AI prompt for deck generation ----
def build_deck_prompt(client_name, industry, problem, geography, audience, style, framework, research_context):
    sections_list = FRAMEWORKS[framework]["sections"]
    sections_str = "\n".join([f"  {i+1}. {s}" for i, s in enumerate(sections_list)])

    research_block = ""
    if research_context.strip():
        research_block = f"""
RESEARCH CONTEXT PROVIDED BY USER:
The user has provided the following research material about the client. Use these specific facts, names, numbers, and quotes in the deck wherever relevant. Do not invent contradictory facts.

{research_context}
"""

    return f"""You are a senior consultant at a top-tier firm ({style} style).
Generate a consulting pitch deck for the following engagement using the {framework} framework.

CLIENT: {client_name}
INDUSTRY: {industry}
BUSINESS PROBLEM: {problem}
GEOGRAPHY: {geography}
AUDIENCE: {audience}
FRAMEWORK: {framework} — {FRAMEWORKS[framework]['description']}

The deck must follow this exact section structure:
{sections_str}
{research_block}
Produce the deck as a JSON object with this exact structure:
{{
  "title_slide": {{
    "title": "<Punchy engagement title>",
    "subtitle": "<One-line description>",
    "client_name": "{client_name}",
    "framework": "{framework}"
  }},
  "slides": [
    {{
      "slide_number": 2,
      "section": "<Section name from the list above>",
      "title": "<Slide title — make it a 'so-what' statement, NOT a topic label>",
      "body_points": ["<3-4 specific points>"],
      "key_takeaway": "<One sentence — the single most important message>",
      "is_options_slide": false
    }}
    ... one slide per section above
  ]
}}

SPECIAL RULE FOR OPTIONS / RECOMMENDATIONS SLIDE:
If the slide presents strategic options (typically the second-to-last slide), set "is_options_slide" to true AND structure body_points as exactly 3 options, formatted as:
  ["Option A — <name>: <2-3 sentence description with pros/cons>", "Option B — <name>: <description>", "Option C — <name>: <description>"]

CRITICAL RULES:
- Every slide title must be a "so-what" statement, not a topic label
  (e.g., "Margins are eroding faster than industry average" NOT "Margin Analysis")
- Body points must be specific — use plausible numbers, named competitors, dates, percentages
- If research context was provided, prioritise those specific facts over invented generic ones
- Key takeaways must be insight-driven, not summary-driven
- Reply with ONLY the JSON object, no markdown code fences, no extra text."""


def call_gemini_for_deck(prompt):
    """Call Gemini and parse the JSON deck structure."""
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
        )
        text = response.text.strip()
        if text.startswith("```"):
            text = text.split("```")[1]
            if text.startswith("json"):
                text = text[4:]
            text = text.strip()
        return json.loads(text)
    except json.JSONDecodeError as e:
        st.error(f"The AI didn't return valid JSON. Try regenerating. Error: {e}")
        return None
    except Exception as e:
        st.error(f"Something went wrong calling the API: {e}")
        return None


# ---- PowerPoint generation ----
def create_pptx(deck_data, palette_name="Corporate Navy"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Pull colours from the chosen palette
    palette = PALETTES.get(palette_name, PALETTES["Corporate Navy"])
    NAVY = palette["primary"]
    GRAY_DARK = palette["text"]
    GRAY_LIGHT = palette["footer"]
    ACCENT = palette["accent"]
    BG_LIGHT = palette["bg_box"]

    blank_layout = prs.slide_layouts[6]
    total_slides = len(deck_data["slides"]) + 1  # +1 for title slide

    # ---------- TITLE SLIDE ----------
    slide = prs.slides.add_slide(blank_layout)

    # Top navy bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.4))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()

    # Diagonal accent shape — gives the title slide visual identity
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(10.5), Inches(0.4),
        Inches(2.83), Inches(2.5)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = ACCENT
    accent_shape.line.fill.background()

    # Bottom accent strip
    bottom_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), prs.slide_width, Inches(0.4))
    bottom_strip.fill.solid()
    bottom_strip.fill.fore_color.rgb = ACCENT
    bottom_strip.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = deck_data["title_slide"]["title"]
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.color.rgb = NAVY

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(1))
    sf = subtitle_box.text_frame
    sf.word_wrap = True
    sf.text = deck_data["title_slide"]["subtitle"]
    sf.paragraphs[0].font.size = Pt(20)
    sf.paragraphs[0].font.name = "Calibri"
    sf.paragraphs[0].font.color.rgb = GRAY_DARK

    # Prepared for
    client_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(12), Inches(0.5))
    cf = client_box.text_frame
    cf.text = f"Prepared for: {deck_data['title_slide']['client_name']}"
    cf.paragraphs[0].font.size = Pt(14)
    cf.paragraphs[0].font.name = "Calibri"
    cf.paragraphs[0].font.color.rgb = ACCENT
    cf.paragraphs[0].font.italic = True

    # Framework used
    fw_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(12), Inches(0.4))
    fwf = fw_box.text_frame
    fwf.text = f"Framework: {deck_data['title_slide'].get('framework', 'Custom')}"
    fwf.paragraphs[0].font.size = Pt(12)
    fwf.paragraphs[0].font.name = "Calibri"
    fwf.paragraphs[0].font.color.rgb = GRAY_LIGHT

    # ---------- CONTENT SLIDES ----------
    for idx, slide_data in enumerate(deck_data["slides"], start=2):
        slide = prs.slides.add_slide(blank_layout)

        # Top navy bar (thin, on every slide)
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = NAVY
        top_bar.line.fill.background()

        # Section label (top-left, small caps, accent colour)
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.4))
        sec_tf = section_box.text_frame
        sec_tf.text = slide_data["section"].upper()
        sec_tf.paragraphs[0].font.size = Pt(11)
        sec_tf.paragraphs[0].font.color.rgb = ACCENT
        sec_tf.paragraphs[0].font.bold = True
        sec_tf.paragraphs[0].font.name = "Calibri"

        # Slide title (the "so-what")
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.text = slide_data["title"]
        title_tf.paragraphs[0].font.size = Pt(26)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = NAVY
        title_tf.paragraphs[0].font.name = "Calibri"

        # Separator line under title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.1), Inches(12.3), Emu(15000))
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY_LIGHT
        line.line.fill.background()

        # ----- Body: options diagram OR bullet points -----
        if slide_data.get("is_options_slide") and len(slide_data["body_points"]) >= 3:
            option_count = min(3, len(slide_data["body_points"]))
            total_width = Inches(12.3)
            gap = Inches(0.3)
            box_width = (total_width - gap * (option_count - 1)) / option_count
            box_height = Inches(3.6)
            box_top = Inches(2.4)

            for i in range(option_count):
                left = Inches(0.5) + (box_width + gap) * i

                bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, box_top, box_width, box_height)
                bg.fill.solid()
                bg.fill.fore_color.rgb = BG_LIGHT
                bg.line.color.rgb = NAVY
                bg.line.width = Pt(1)

                header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, box_top, box_width, Inches(0.5))
                header.fill.solid()
                header.fill.fore_color.rgb = NAVY
                header.line.fill.background()

                option_label = f"OPTION {chr(65 + i)}"
                header_text = slide.shapes.add_textbox(left, box_top, box_width, Inches(0.5))
                ht = header_text.text_frame
                ht.text = option_label
                ht.paragraphs[0].font.size = Pt(14)
                ht.paragraphs[0].font.bold = True
                ht.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                ht.paragraphs[0].font.name = "Calibri"
                ht.paragraphs[0].alignment = 2

                body_text = slide.shapes.add_textbox(left + Inches(0.15), box_top + Inches(0.6), box_width - Inches(0.3), box_height - Inches(0.7))
                bt = body_text.text_frame
                bt.word_wrap = True
                bt.text = slide_data["body_points"][i]
                bt.paragraphs[0].font.size = Pt(12)
                bt.paragraphs[0].font.color.rgb = GRAY_DARK
                bt.paragraphs[0].font.name = "Calibri"

        else:
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.7))
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for i, point in enumerate(slide_data["body_points"]):
                if i == 0:
                    p = body_tf.paragraphs[0]
                else:
                    p = body_tf.add_paragraph()
                p.text = f"•  {point}"
                p.font.size = Pt(16)
                p.font.color.rgb = GRAY_DARK
                p.font.name = "Calibri"
                p.space_after = Pt(14)

        # Key takeaway bar
        takeaway_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7))
        takeaway_bar.fill.solid()
        takeaway_bar.fill.fore_color.rgb = NAVY
        takeaway_bar.line.fill.background()

        takeaway_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.6))
        tk_tf = takeaway_text.text_frame
        tk_tf.word_wrap = True
        p = tk_tf.paragraphs[0]
        p.text = f"KEY TAKEAWAY:  {slide_data['key_takeaway']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.name = "Calibri"

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3))
        ft = footer_box.text_frame
        ft.text = f"{deck_data['title_slide']['client_name']}    |    Strictly confidential"
        ft.paragraphs[0].font.size = Pt(9)
        ft.paragraphs[0].font.color.rgb = GRAY_LIGHT
        ft.paragraphs[0].font.name = "Calibri"

        # Page number
        page_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3))
        pt = page_box.text_frame
        pt.text = f"{idx} / {total_slides}"
        pt.paragraphs[0].font.size = Pt(9)
        pt.paragraphs[0].font.color.rgb = GRAY_LIGHT
        pt.paragraphs[0].font.name = "Calibri"
        pt.paragraphs[0].alignment = 2

    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes


# ---- User interface ----
st.markdown("Describe a consulting engagement, pick a framework, and generate a structured pitch deck.")

col1, col2 = st.columns(2)
with col1:
    client_name = st.text_input("Client name", placeholder="AirAsia")
    industry = st.text_input("Industry", placeholder="Aviation / Low-cost carriers")
    geography = st.text_input("Geographic focus", placeholder="Southeast Asia")
with col2:
    audience = st.text_input("Target audience", placeholder="CEO and senior leadership team")
    style = st.selectbox("Style", ["McKinsey", "BCG", "Big Four (Deloitte/PwC/EY/KPMG)"])
    framework = st.selectbox(
        "Analytical framework",
        list(FRAMEWORKS.keys()),
        help="Different frameworks fit different problem types",
    )
    suggested = suggest_palette(client_name or "", industry or "")
    palette_options = list(PALETTES.keys())
    default_index = palette_options.index(suggested) if suggested in palette_options else 0
    palette_choice = st.selectbox(
        "Visual theme",
        palette_options,
        index=default_index,
        help=f"Auto-suggested based on industry: {suggested}",
    )

st.caption(f"📖 **{framework}** — {FRAMEWORKS[framework]['description']}")

problem = st.text_area(
    "Business problem",
    placeholder="Declining domestic passenger volumes as competitors expand low-cost routes. Margin pressure from rising fuel costs.",
    height=100,
)

with st.expander("➕ Add research context (optional but recommended)", expanded=False):
    st.caption("Paste real research about the client — recent press releases, earnings call notes, industry data. The AI will use these specific facts instead of generic ones.")
    research_context = st.text_area(
        "Research context",
        placeholder="""Example:
- AirAsia reported Q3 2025 load factor of 87%, up from 82% YoY
- New CEO Tony Fernandes announced 'Capital A 2.0' strategy in October 2025
- Main competitors: Scoot (Singapore), VietJet (Vietnam), Cebu Pacific (Philippines)
- Recent challenge: fuel costs up 18% in 2025, hedging contracts expiring Q2 2026""",
        height=180,
        label_visibility="collapsed",
    )

if st.button("Generate pitch deck", type="primary"):
    if not all([client_name, industry, problem, geography, audience]):
        st.warning("Please fill in all required fields before generating.")
    else:
        with st.spinner(f"Drafting the deck using {framework}... (10-25 seconds)"):
            prompt = build_deck_prompt(client_name, industry, problem, geography, audience, style, framework, research_context)
            deck_data = call_gemini_for_deck(prompt)

            if deck_data:
                st.success(f"Deck generated using {framework}! Preview below — download to open in PowerPoint.")

                with st.expander("Preview deck structure", expanded=True):
                    st.subheader(deck_data["title_slide"]["title"])
                    st.caption(deck_data["title_slide"]["subtitle"])
                    st.caption(f"Framework: {deck_data['title_slide'].get('framework', framework)}")
                    st.divider()
                    for s in deck_data["slides"]:
                        st.markdown(f"**Slide {s['slide_number']} — {s['section']}**")
                        st.markdown(f"### {s['title']}")
                        for p in s['body_points']:
                            st.markdown(f"- {p}")
                        st.info(f"💡 {s['key_takeaway']}")
                        st.divider()

                pptx_file = create_pptx(deck_data, palette_name=palette_choice)
                filename = f"{client_name.replace(' ', '_')}_{framework.split()[0]}_deck.pptx"
                st.download_button(
                    label="📥 Download .pptx file",
                    data=pptx_file,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

with st.sidebar:
    st.header("About this tool")
    st.markdown("""
    Generates structured consulting pitch decks based on a client situation.
    
    **What's new in this version:**
    - Research context field for grounding output in real client facts
    - 6 analytical frameworks (SCR, Porter's Five Forces, SWOT, McKinsey 7S, Three Horizons, MECE)
    - 8 visual themes auto-matched to industry
    - Polished PowerPoint design with header strips, page numbers, footer
    - Strategic options shown as a visual 3-column comparison

    **What it doesn't do:**
    - Doesn't know your real client
    - Can't replace consultant judgment on specifics
    - First-draft accelerator, not final-draft tool
    """)
    st.divider()
    st.caption("Built by [Gan Kai Li](https://www.linkedin.com/in/gan-kai-li-a8a782317/) · ICAEW student preparing for Big Four consulting")